import os
import random
import google.generativeai as genai
from flask import Flask, render_template, request, jsonify, send_file
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import telegram

app = Flask(__name__)

# إعداد Gemini
genai.configure(api_key="AIzaSyAYQeOd-6H1xHxOhFxM2J4r7p4WHsBPPUg")
model = genai.GenerativeModel('gemini-1.5-pro')

# إعداد بوت Telegram
TELEGRAM_BOT_TOKEN = "8214786867:AAHsLBghSsF2le7Tx_rsLQd6GaXFWVgs_GA"
CHAT_ID = "7836619198"  # يمكن يكون رقم القناة أو المستخدم
bot = telegram.Bot(token=TELEGRAM_BOT_TOKEN)

def generate_unique_legal_letter(data):
    styles = [
        "أسلوب دستوري قانوني رصين",
        "أسلوب إداري حازم",
        "أسلوب بلاغي سيادي فريد",
        "أسلوب رقابي مشدد"
    ]
    selected_style = random.choice(styles)
    
    prompt = f"""
    بصفتك خبيراً في صياغة كتب مجلس النواب العراقي، صغ كتاباً رسمياً من (النائب حيدر الأسدي).
    الموجه إلى: {data['target_agency']}
    الموضوع: {data['subject']}
    التفاصيل: {data['details']}
    
    الضوابط: صياغة من الصفر، {selected_style}، لغة قانونية فصحى، بدون قوالب، بدون "تحية طيبة".
    """
    response = model.generate_content(prompt)
    return response.text

def save_to_pptx(text, filename):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.RIGHT
    prs.save(filename)

def send_to_telegram(file_path):
    bot.send_document(chat_id=CHAT_ID, document=open(file_path, "rb"))

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/process', methods=['POST'])
def process():
    data = request.json
    legal_text = generate_unique_legal_letter(data)
    filename = f"letter_{random.randint(1000,9999)}.pptx"
    save_to_pptx(legal_text, filename)
    
    # إرسال الملف مباشرة عبر بوت Telegram
    send_to_telegram(filename)
    
    return jsonify({"text": legal_text, "file": filename})

@app.route('/download/<filename>')
def download(filename):
    return send_file(filename, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
